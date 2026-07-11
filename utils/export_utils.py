from io import BytesIO
from docx import Document
from pptx import Presentation

def generate_word_report(data):

    doc = Document()

    doc.add_heading("Enterprise AI Transformation Report", level=1)

    company = data["company_data"]["company_profile"]

    doc.add_heading("Company", level=2)

    for k, v in company.items():
        doc.add_paragraph(f"{k}: {v}")

    roi = data["roi_results"]

    doc.add_heading("ROI", level=2)

    doc.add_paragraph(f"Investment : ${roi['investment']:,.0f}")
    doc.add_paragraph(f"Benefit : ${roi['benefit']:,.0f}")
    doc.add_paragraph(f"ROI : {roi['roi']:.1f}%")
    doc.add_paragraph(f"Payback : {roi['payback']:.2f} Years")

    knowledge = data["knowledge_assessment"]

    doc.add_heading("Knowledge Assessment", level=2)

    doc.add_paragraph(
        f"Overall Score : {knowledge['overall_score']}%"
    )

    buffer = BytesIO()

    doc.save(buffer)

    buffer.seek(0)

    return buffer


def generate_ppt(data):

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])

    slide.shapes.title.text = "Enterprise AI Transformation"

    slide.placeholders[1].text = "Executive Report"

    slide = prs.slides.add_slide(prs.slide_layouts[5])

    slide.shapes.title.text = "ROI Summary"

    roi = data["roi_results"]

    textbox = slide.shapes.add_textbox(
        1000000,
        1000000,
        7000000,
        3000000
    )

    tf = textbox.text_frame

    tf.text = f"Investment : ${roi['investment']:,.0f}"

    tf.add_paragraph().text = f"Benefit : ${roi['benefit']:,.0f}"

    tf.add_paragraph().text = f"ROI : {roi['roi']:.1f}%"

    tf.add_paragraph().text = f"Payback : {roi['payback']:.2f} Years"

    buffer = BytesIO()

    prs.save(buffer)

    buffer.seek(0)

    return buffer